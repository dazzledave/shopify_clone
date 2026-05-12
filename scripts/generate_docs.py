import os
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

def create_pptx(content, output_file):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Shopify Clone: Technical Documentation"
    subtitle.text = "Project Overview and Architecture\nCreated by Antigravity"

    sections = content.split("## ")[1:]
    for section in sections:
        lines = section.split("\n")
        section_title = lines[0]
        section_content = "\n".join(lines[1:]).strip()
        
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = section_title
        tf = body_shape.text_frame
        tf.text = section_content[:500] + "..." if len(section_content) > 500 else section_content

    prs.save(output_file)
    print(f"PPTX saved to {output_file}")

def create_pdf(content, output_file):
    doc = SimpleDocTemplate(output_file, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    lines = content.split("\n")
    for line in lines:
        if line.startswith("# "):
            story.append(Paragraph(line[2:], styles['Title']))
        elif line.startswith("## "):
            story.append(Paragraph(line[3:], styles['Heading2']))
        elif line.startswith("### "):
            story.append(Paragraph(line[4:], styles['Heading3']))
        elif line.strip():
            story.append(Paragraph(line, styles['Normal']))
        story.append(Spacer(1, 12))

    doc.build(story)
    print(f"PDF saved to {output_file}")

if __name__ == "__main__":
    with open("full_documentation.md", "r", encoding="utf-8") as f:
        md_content = f.read()
    
    create_pptx(md_content, "shopify_clone_docs.pptx")
    create_pdf(md_content, "shopify_clone_docs.pdf")
